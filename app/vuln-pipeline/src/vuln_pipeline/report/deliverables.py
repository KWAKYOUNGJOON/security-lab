from __future__ import annotations

from pathlib import Path
from typing import Any

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.shared import Inches, Pt

from vuln_pipeline.models import ReportBundle
from vuln_pipeline.report.docx import DocxRenderer
from vuln_pipeline.report.knowledge import load_deliverable_profile
from vuln_pipeline.report.markdown import render_markdown_report
from vuln_pipeline.storage import write_csv, write_json, write_markdown
from vuln_pipeline.utils import ensure_directory, now_utc


def generate_deliverables(
    bundle: ReportBundle,
    contexts: dict[str, dict[str, Any]],
    review_queue: list[dict[str, Any]],
    review_closure_status: dict[str, Any],
    readiness: dict[str, Any],
    deliverables_root: Path,
    deliverable_profile_dir: Path,
    deliverable_profile_name: str,
    generate_docx: bool,
    final_delivery_included: bool = False,
    real_input_selection: dict[str, Any] | None = None,
) -> dict[str, Any]:
    profile = load_deliverable_profile(deliverable_profile_dir, deliverable_profile_name)
    ensure_directory(deliverables_root)
    customer_meta = contexts["customer"]["document_meta"]
    suffix = f"_{contexts['customer']['profile']['id']}_{customer_meta.get('report_version', 'v1.0')}" if profile.get("versioned_filenames") else ""
    included_files: list[str] = []
    excluded_files: list[str] = []
    customer_outputs: dict[str, list[str]] = {}
    internal_outputs: dict[str, list[str]] = {}
    closeout_summary = _closeout_summary(review_closure_status)
    presentation_result: dict[str, Any] = {"status": "skipped", "pptx_generated": False}

    if profile.get("full_reports", True):
        for audience, base_name in [("internal", "full_report_internal"), ("customer", "full_report_customer")]:
            context = contexts[audience]
            output_list = internal_outputs if audience == "internal" else customer_outputs
            md_path = deliverables_root / f"{base_name}{suffix}.md"
            write_markdown(md_path, render_markdown_report(context))
            included_files.append(str(md_path))
            output_list.setdefault(f"{audience}_full_report", []).append(str(md_path))
            if generate_docx:
                docx_path = deliverables_root / f"{base_name}{suffix}.docx"
                DocxRenderer().render(context, docx_path)
                included_files.append(str(docx_path))
                output_list.setdefault(f"{audience}_full_report", []).append(str(docx_path))

    if profile.get("onepager", True):
        audience = profile.get("onepager_audience", "customer")
        onepager = build_executive_onepager(contexts[audience], review_queue, review_closure_status, readiness, audience)
        md_path = deliverables_root / f"executive_onepager{suffix}.md"
        write_markdown(md_path, render_onepager_markdown(onepager))
        included_files.append(str(md_path))
        if audience == "customer":
            customer_outputs.setdefault("customer_onepager", []).append(str(md_path))
        else:
            internal_outputs.setdefault("internal_onepager", []).append(str(md_path))
        if generate_docx:
            docx_path = deliverables_root / f"executive_onepager{suffix}.docx"
            render_onepager_docx(onepager, docx_path)
            included_files.append(str(docx_path))
            if audience == "customer":
                customer_outputs.setdefault("customer_onepager", []).append(str(docx_path))
            else:
                internal_outputs.setdefault("internal_onepager", []).append(str(docx_path))

    tracker_rows: list[dict[str, Any]] = []
    if profile.get("tracker", True):
        tracker_rows = build_remediation_tracker(bundle, review_queue, bundle.comparison_summary, final_delivery_included, review_closure_status)
        csv_path = deliverables_root / f"remediation_tracker{suffix}.csv"
        md_path = deliverables_root / f"remediation_tracker{suffix}.md"
        write_csv(csv_path, tracker_rows, list(tracker_rows[0].keys()) if tracker_rows else _tracker_columns())
        write_markdown(md_path, render_tracker_markdown(tracker_rows))
        included_files.extend([str(csv_path), str(md_path)])
        customer_outputs.setdefault("customer_tracker", []).append(str(csv_path))
        customer_outputs.setdefault("customer_tracker", []).append(str(md_path))

    if profile.get("handoff", True):
        handoff = build_analyst_handoff(bundle, contexts["internal"], review_queue, review_closure_status, readiness)
        md_path = deliverables_root / f"analyst_handoff{suffix}.md"
        write_markdown(md_path, handoff)
        included_files.append(str(md_path))
        internal_outputs.setdefault("analyst_handoff", []).append(str(md_path))
    else:
        excluded_files.append("analyst_handoff")

    if profile.get("presentation", True):
        audience = profile.get("presentation_audience", "customer")
        presentation = build_presentation_data(bundle, contexts[audience], review_queue, review_closure_status, readiness)
        outline_path = deliverables_root / f"presentation_outline{suffix}.md"
        data_path = deliverables_root / f"presentation_data{suffix}.json"
        write_markdown(outline_path, render_presentation_outline(presentation))
        write_json(data_path, presentation)
        included_files.extend([str(outline_path), str(data_path)])
        internal_outputs.setdefault("presentation_support", []).extend([str(outline_path), str(data_path)])
        pptx_path = deliverables_root / f"presentation_briefing{suffix}.pptx"
        presentation_result = render_presentation_pptx(presentation, pptx_path, data_path)
        if presentation_result["status"] == "pptx":
            included_files.append(str(pptx_path))
            customer_outputs.setdefault("customer_presentation", []).append(str(pptx_path))
            internal_outputs.setdefault("presentation_support", []).append(str(pptx_path))
        else:
            fallback_path = deliverables_root / f"presentation_briefing{suffix}_fallback.json"
            write_json(fallback_path, presentation_result)
            included_files.append(str(fallback_path))
            internal_outputs.setdefault("presentation_support", []).append(str(fallback_path))
            if profile.get("customer_submission_include_presentation_fallback", False):
                customer_outputs.setdefault("customer_presentation", []).append(str(fallback_path))
        presentation_result["outline_path"] = str(outline_path)
        presentation_result["data_path"] = str(data_path)

    checklist = build_review_closure_checklist(bundle, review_queue, review_closure_status, readiness)
    checklist_path = deliverables_root / f"review_closure_checklist{suffix}.md"
    write_markdown(checklist_path, checklist)
    included_files.append(str(checklist_path))
    internal_outputs.setdefault("review_closure_checklist", []).append(str(checklist_path))

    submission_memo = build_submission_memo(
        bundle=bundle,
        customer_context=contexts["customer"],
        review_closure_status=review_closure_status,
        readiness=readiness,
        included_customer_files=_dedupe_paths(
            customer_outputs.get("customer_full_report", [])
            + customer_outputs.get("customer_onepager", [])
            + customer_outputs.get("customer_tracker", [])
            + customer_outputs.get("customer_presentation", [])
        ),
    )
    submission_memo_path = deliverables_root / "submission_memo.md"
    write_markdown(submission_memo_path, render_submission_memo_markdown(submission_memo))
    included_files.append(str(submission_memo_path))
    internal_outputs.setdefault("submission_memo", []).append(str(submission_memo_path))

    real_data_summary = build_real_data_rehearsal_summary(
        bundle=bundle,
        review_queue=review_queue,
        readiness=readiness,
        real_input_selection=real_input_selection,
    )
    real_data_summary_path = deliverables_root / "real_data_rehearsal_summary.md"
    write_markdown(real_data_summary_path, render_real_data_rehearsal_summary(real_data_summary))
    included_files.append(str(real_data_summary_path))
    internal_outputs.setdefault("real_data_rehearsal_summary", []).append(str(real_data_summary_path))

    return {
        "run_id": bundle.run_id,
        "report_version": customer_meta.get("report_version", "v1.0"),
        "delivery_date": customer_meta.get("delivery_date", customer_meta.get("generated_at", "")),
        "deliverable_profile": profile["id"],
        "report_profile": bundle.report_profile,
        "report_template": bundle.report_template,
        "readiness_policy": readiness.get("policy_name"),
        "included_files": _dedupe_paths(included_files),
        "excluded_files": excluded_files,
        "customer_outputs": customer_outputs,
        "internal_outputs": internal_outputs,
        "packaging_time": now_utc(),
        "checksums_ref": str(Path("delivery") / "checksums.json"),
        "presentation_status": presentation_result["status"],
        "pptx_generated": bool(presentation_result.get("pptx_generated")),
        "pptx_validation": presentation_result,
        "final_delivery_included": final_delivery_included,
        "closeout_summary": closeout_summary,
        "customer_redaction_applied": not contexts["customer"]["profile"].get("show_internal_paths", True),
        "branding_applied": bool(customer_meta.get("branding_applied")),
        "submission_memo": str(submission_memo_path),
        "real_data_rehearsal_summary": str(real_data_summary_path),
    }


def build_executive_onepager(
    context: dict[str, Any],
    review_queue: list[dict[str, Any]],
    review_closure_status: dict[str, Any],
    readiness: dict[str, Any],
    audience: str,
) -> dict[str, Any]:
    bundle = context["bundle"]
    meta = context["document_meta"]
    top_issues = [item["issue"] for item in context["narratives"][:3]]
    comparison = context["comparison_summary"]
    high_count = context["severity_stats"].get("Critical", 0) + context["severity_stats"].get("High", 0)
    top_review = review_queue[:3]
    remaining_blocker_count = review_closure_status.get("unresolved_review_items", 0)
    return {
        "title": meta.get("cover_title") or "Executive One-Pager",
        "subtitle": meta.get("subtitle") or ("Customer delivery summary" if audience == "customer" else "Internal release summary"),
        "meta": meta,
        "key_conclusions": [
            f"Total issues: {bundle.summary.get('issues', 0)}. High or above: {high_count}.",
            f"Resolved review items: {review_closure_status.get('resolved_review_items', 0)}. Accepted risk items: {review_closure_status.get('accepted_risk_items', 0)}.",
            f"Baseline comparison new/resolved/changed/unchanged: {len(comparison.get('new_issues', []))}/{len(comparison.get('resolved_issues', []))}/{len(comparison.get('changed_issues', []))}/{len(comparison.get('unchanged_issues', []))}.",
        ],
        "counts": {
            "total_issues": bundle.summary.get("issues", 0),
            "high_or_above": high_count,
            "immediate_action_required": sum(1 for item in context["narratives"] if item["issue"].severity.level in {"Critical", "High"}),
            "new_issues": len(comparison.get("new_issues", [])),
            "resolved_issues": len(comparison.get("resolved_issues", [])),
            "changed_issues": len(comparison.get("changed_issues", [])),
            "unchanged_issues": len(comparison.get("unchanged_issues", [])),
            "unresolved_review": review_closure_status.get("unresolved_review_items", 0),
            "accepted_risk": review_closure_status.get("accepted_risk_items", 0),
            "resolved_review": review_closure_status.get("resolved_review_items", 0),
            "remaining_blocker": remaining_blocker_count,
        },
        "top_risks": [
            {"issue_id": issue.issue_id, "title": issue.title, "severity": issue.severity.level, "assets": issue.affected_assets}
            for issue in top_issues
        ],
        "asset_summary": sorted(context["asset_groups"].keys()),
        "priorities": context["action_plan"],
        "top_review_targets": top_review,
        "review_status": {
            "analyst_review_required": review_closure_status.get("unresolved_review_items", 0) > 0,
            "suppression_count": len(bundle.suppressed_issues),
            "comparison_available": comparison.get("available", False),
            "release_readiness": readiness["status"],
        },
        "readiness": readiness,
    }


def render_onepager_markdown(onepager: dict[str, Any]) -> str:
    meta = onepager["meta"]
    lines = [
        f"# {onepager['title']}",
        "",
        f"> {onepager['subtitle']}",
        "",
        f"- client_name: `{meta.get('client_name', '')}`",
        f"- project_name: `{meta.get('project_name', '')}`",
        f"- report_version: `{meta.get('report_version', '')}`",
        f"- approver_name: `{meta.get('approver_name', '')}`",
        f"- approver_title: `{meta.get('approver_title', '')}`",
        f"- delivery_date: `{meta.get('delivery_date', meta.get('generated_at', ''))}`",
        f"- contact_email: `{meta.get('contact_email', '')}`",
        f"- run_id: `{meta.get('run_id', '')}`",
        "",
        "## Key Conclusions",
    ]
    lines.extend([f"- {item}" for item in onepager["key_conclusions"]])
    lines.extend(
        [
            "",
            "## Key Counts",
            f"- total issues: {onepager['counts']['total_issues']}",
            f"- high or above: {onepager['counts']['high_or_above']}",
            f"- immediate action required: {onepager['counts']['immediate_action_required']}",
            f"- resolved review count: {onepager['counts']['resolved_review']}",
            f"- accepted risk count: {onepager['counts']['accepted_risk']}",
            f"- remaining blocker count: {onepager['counts']['remaining_blocker']}",
            f"- new/resolved/changed/unchanged: {onepager['counts']['new_issues']}/{onepager['counts']['resolved_issues']}/{onepager['counts']['changed_issues']}/{onepager['counts']['unchanged_issues']}",
            "",
            "## Top Risks",
        ]
    )
    for item in onepager["top_risks"]:
        lines.append(f"- {item['issue_id']} {item['title']} ({item['severity']}) / {', '.join(item['assets'])}")
    lines.extend(["", "## Review Priority"])
    for item in onepager["top_review_targets"] or [{"issue_id": "-", "priority_band": "-", "recommended_action": "none"}]:
        lines.append(f"- {item['issue_id']} ({item.get('priority_band', '-')}, {item.get('suggested_sla', '-')}) {item['recommended_action']}")
    lines.extend(["", "## Readiness"])
    lines.append(f"- status: {onepager['review_status']['release_readiness']}")
    lines.append(f"- analyst review required: {onepager['review_status']['analyst_review_required']}")
    if meta.get("footer_notice"):
        lines.extend(["", "## Notice", f"- {meta['footer_notice']}"])
    return "\n".join(lines) + "\n"


def render_onepager_docx(onepager: dict[str, Any], output_path: Path) -> None:
    document = _base_document(onepager["title"])
    meta = onepager["meta"]
    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run(onepager["title"])
    run.bold = True
    run.font.size = Pt(22)
    subtitle = document.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.add_run(onepager["subtitle"])
    branding = document.add_paragraph()
    branding.alignment = WD_ALIGN_PARAGRAPH.CENTER
    branding.add_run(" | ".join(item for item in [meta.get("organization_name", ""), meta.get("client_name", ""), meta.get("project_name", "")] if item))
    if meta.get("logo_path_optional"):
        logo = document.add_paragraph()
        logo.alignment = WD_ALIGN_PARAGRAPH.CENTER
        logo.add_run(f"Logo: {Path(meta['logo_path_optional']).name if Path(meta['logo_path_optional']).exists() else 'not provided'}")
    table = document.add_table(rows=7, cols=2)
    table.style = "Table Grid"
    for idx, row in enumerate(
        [
            ("Client", meta.get("client_name", "")),
            ("Project", meta.get("project_name", "")),
            ("Version", meta.get("report_version", "")),
            ("Approver", meta.get("approver_name", "")),
            ("Approver Title", meta.get("approver_title", "")),
            ("Delivery Date", meta.get("delivery_date", meta.get("generated_at", ""))),
            ("Contact", meta.get("contact_email", "")),
        ]
    ):
        table.rows[idx].cells[0].text = row[0]
        table.rows[idx].cells[1].text = str(row[1])
    signature = document.add_table(rows=2, cols=2)
    signature.style = "Table Grid"
    signature.rows[0].cells[0].text = "Prepared By"
    signature.rows[0].cells[1].text = meta.get("analyst_name", "")
    signature.rows[1].cells[0].text = "Approved By"
    signature.rows[1].cells[1].text = " / ".join(item for item in [meta.get("approver_name", ""), meta.get("approver_title", "")] if item)
    document.add_heading("Key Conclusions", level=1)
    for item in onepager["key_conclusions"]:
        document.add_paragraph(item, style="List Bullet")
    document.add_heading("Review Priority", level=1)
    for item in onepager["top_review_targets"] or [{"issue_id": "-", "priority_band": "-", "recommended_action": "none"}]:
        document.add_paragraph(f"{item['issue_id']} {item.get('priority_band', '-')} {item['recommended_action']}", style="List Bullet")
    if meta.get("footer_notice"):
        document.add_heading("Notice", level=1)
        document.add_paragraph(meta["footer_notice"])
    output_path.parent.mkdir(parents=True, exist_ok=True)
    document.save(output_path)


def build_remediation_tracker(
    bundle: ReportBundle,
    review_queue: list[dict[str, Any]],
    comparison: dict[str, Any],
    final_delivery_included: bool,
    review_closure_status: dict[str, Any],
) -> list[dict[str, Any]]:
    review_map = {row["issue_id"]: row for row in review_queue}
    compare_map = _comparison_status_map(comparison)
    rows: list[dict[str, Any]] = []
    closeout_status = "closed" if review_closure_status.get("unresolved_review_items", 0) == 0 else "open"
    for issue in list(bundle.issues) + list(bundle.suppressed_issues):
        review = review_map.get(issue.issue_id, {})
        rows.append(
            {
                "issue_id": issue.issue_id,
                "title": issue.title,
                "severity": issue.severity.level,
                "confidence": issue.confidence.level,
                "weakness_family": issue.weakness_family or "",
                "primary_cwe": issue.primary_cwe or "",
                "affected_assets": "; ".join(issue.affected_assets),
                "immediate_action": _first(issue.remediation, "Review report recommendations"),
                "structural_fix": issue.remediation[1] if len(issue.remediation) > 1 else _first(issue.remediation, "Apply structural control improvements"),
                "validation_after_fix": issue.evidence_summary[0] if issue.evidence_summary else "Retest the affected endpoint after remediation.",
                "recommended_owner": issue.recommended_owner or "AppSec",
                "target_due": issue.target_due or "TBD",
                "status": issue.suppression_status or ("false_positive" if issue.false_positive else "open"),
                "source_run_id": bundle.run_id,
                "compared_run_status": compare_map.get(issue.cluster_key, "n/a"),
                "suppression_status": issue.suppression_status or "",
                "review_required": "yes" if issue.issue_id in review_map and not review.get("is_resolved") else "no",
                "review_disposition": review.get("review_disposition", ""),
                "approval_status": review.get("approval_status", "pending"),
                "final_delivery_included": "yes" if final_delivery_included else "no",
                "closeout_status": closeout_status,
            }
        )
    return rows


def render_tracker_markdown(rows: list[dict[str, Any]]) -> str:
    lines = [
        "# Remediation Tracker",
        "",
        "| Issue ID | Severity | Owner | Target Due | Compare Status | Review | Disposition | Approval | Final | Closeout |",
        "|---|---|---|---|---|---|---|---|---|---|",
    ]
    for row in rows:
        lines.append(
            f"| {row['issue_id']} | {row['severity']} | {row['recommended_owner']} | {row['target_due']} | "
            f"{row['compared_run_status']} | {row['review_required']} | {row['review_disposition'] or '-'} | "
            f"{row['approval_status']} | {row['final_delivery_included']} | {row['closeout_status']} |"
        )
    if not rows:
        lines.append("| - | - | - | - | - | - | - | - | - | - |")
    return "\n".join(lines) + "\n"


def build_analyst_handoff(
    bundle: ReportBundle,
    internal_context: dict[str, Any],
    review_queue: list[dict[str, Any]],
    review_closure_status: dict[str, Any],
    readiness: dict[str, Any],
) -> str:
    lines = [
        "# Analyst Handoff",
        "",
        "## Execution",
        f"- run_id: {bundle.run_id}",
        f"- generated_at: {bundle.generated_at}",
        f"- issues: {bundle.summary.get('issues', 0)}",
        "",
        "## Closeout Summary",
        f"- unresolved review items: {review_closure_status.get('unresolved_review_items', 0)}",
        f"- accepted risk items: {review_closure_status.get('accepted_risk_items', 0)}",
        f"- deferred items: {review_closure_status.get('deferred_items', 0)}",
        f"- readiness: {readiness['status']}",
        "",
        "## Review Queue Top Items",
    ]
    for row in review_queue[:5] or [{"issue_id": "-", "recommended_action": "none"}]:
        lines.append(f"- {row['issue_id']}: {row['recommended_action']}")
    lines.extend(["", "## Next Actions"])
    for section, items in internal_context["action_plan"].items():
        lines.append(f"### {section}")
        for item in items or ["none"]:
            lines.append(f"- {item}")
    return "\n".join(lines) + "\n"


def build_presentation_data(
    bundle: ReportBundle,
    context: dict[str, Any],
    review_queue: list[dict[str, Any]],
    review_closure_status: dict[str, Any],
    readiness: dict[str, Any],
) -> dict[str, Any]:
    comparison = bundle.comparison_summary
    meta = context["document_meta"]
    return {
        "run_id": bundle.run_id,
        "meta": meta,
        "readiness": readiness,
        "slides": [
            {
                "title": meta.get("cover_title") or "Presentation Briefing",
                "summary": meta.get("subtitle") or f"{meta.get('project_name')} / {meta.get('client_name')}",
            },
            {"title": "Project Overview", "summary": context["executive_summary"]},
            {
                "title": "Scope and Inputs",
                "summary": f"Input files: {len(bundle.input_files)} / Affected assets: {len(context['asset_groups'])}",
            },
            {
                "title": "Assessment Summary",
                "summary": f"Issues: {bundle.summary.get('issues', 0)} / Unresolved review: {review_closure_status.get('unresolved_review_items', 0)}",
            },
            {
                "title": "Top Risks",
                "items": [
                    {"issue_id": item["issue"].issue_id, "title": item["issue"].title, "severity": item["issue"].severity.level}
                    for item in context["narratives"][:3]
                ],
            },
            {"title": "Action Plan", "summary": context["action_plan"]},
            {
                "title": "Baseline Comparison",
                "summary": {
                    "new": len(comparison.get("new_issues", [])),
                    "resolved": len(comparison.get("resolved_issues", [])),
                    "changed": len(comparison.get("changed_issues", [])),
                    "unchanged": len(comparison.get("unchanged_issues", [])),
                },
            },
            {"title": "Readiness and Next Steps", "summary": readiness["summary"]["recommendation"]},
        ],
    }


def render_presentation_outline(presentation: dict[str, Any]) -> str:
    lines = ["# Presentation Outline", ""]
    for slide in presentation["slides"]:
        lines.append(f"## {slide['title']}")
        if "summary" in slide:
            lines.append(f"- {slide['summary']}")
        for item in slide.get("items", []):
            lines.append(f"- {item['issue_id']} {item['title']} ({item['severity']})")
        lines.append("")
    return "\n".join(lines)


def render_presentation_pptx(presentation: dict[str, Any], output_path: Path, source_data_path: Path) -> dict[str, Any]:
    expected_slide_count = len(presentation["slides"])
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as exc:
        return {
            "status": "fallback",
            "pptx_generated": False,
            "reason": f"missing dependency: {exc}",
            "missing_dependency_reason": str(exc),
            "install_hint": "pip install python-pptx",
            "expected_output": str(output_path),
            "source_presentation_data": str(source_data_path),
            "expected_minimum_slide_count": expected_slide_count,
        }
    try:
        prs = Presentation()
        for index, slide_data in enumerate(presentation["slides"], start=1):
            layout = prs.slide_layouts[0] if index == 1 else prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_data["title"]
            body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
            summary_text = str(slide_data.get("summary", ""))
            if body is not None:
                body.text = summary_text
                for item in slide_data.get("items", []):
                    body.text_frame.add_paragraph().text = f"{item['issue_id']} {item['title']} ({item['severity']})"
            else:
                textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
                textbox.text_frame.text = summary_text
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)

        saved = Presentation(output_path)
        titles = [slide.shapes.title.text.strip() for slide in saved.slides if slide.shapes.title and slide.shapes.title.text]
        title_check = bool(titles and titles[0] == presentation["slides"][0]["title"])
        text_check = any(item["title"] in titles for item in presentation["slides"][1:]) if presentation["slides"][1:] else True
        slide_count_check = len(saved.slides) >= expected_slide_count
        status = "pptx" if title_check and slide_count_check else "fallback"
        result = {
            "status": status,
            "pptx_generated": status == "pptx",
            "output": str(output_path),
            "source_presentation_data": str(source_data_path),
            "slide_count": len(saved.slides),
            "expected_minimum_slide_count": expected_slide_count,
            "slide_count_valid": slide_count_check,
            "title_validation_passed": title_check,
            "core_text_validation_passed": text_check,
        }
        if status != "pptx":
            result["reason"] = "pptx validation failed after generation"
            result["install_hint"] = "Verify generated file contents and python-pptx compatibility."
            result["expected_output"] = str(output_path)
        return result
    except Exception as exc:
        return {
            "status": "fallback",
            "pptx_generated": False,
            "reason": str(exc),
            "missing_dependency_reason": "",
            "install_hint": "Verify python-pptx installation and template compatibility.",
            "expected_output": str(output_path),
            "source_presentation_data": str(source_data_path),
            "expected_minimum_slide_count": expected_slide_count,
        }


def build_review_closure_checklist(
    bundle: ReportBundle,
    review_queue: list[dict[str, Any]],
    review_closure_status: dict[str, Any],
    readiness: dict[str, Any],
) -> str:
    lines = ["# Review Closure Checklist", ""]
    checks = [
        ("Review queue total/unresolved", f"{review_closure_status.get('total_review_items', 0)}/{review_closure_status.get('unresolved_review_items', 0)}", "DONE" if review_closure_status.get("unresolved_review_items", 0) == 0 else "PENDING"),
        ("Remaining P1/P2", f"{review_closure_status.get('remaining_p1', 0)}/{review_closure_status.get('remaining_p2', 0)}", "DONE" if review_closure_status.get("remaining_p1", 0) == 0 and review_closure_status.get("remaining_p2", 0) == 0 else "BLOCKED"),
        ("Override applied", str(bool(bundle.override_summary.get("applied"))), "DONE" if bundle.override_summary.get("applied") else "PENDING"),
        ("Suppression approved", str(bool(bundle.suppressed_issues)), "DONE" if bundle.suppressed_issues else "PENDING"),
        ("High severity handled", "true" if not any(issue.severity.level in {"Critical", "High"} for issue in bundle.issues) else "false", "BLOCKED" if any(issue.severity.level in {"Critical", "High"} for issue in bundle.issues) else "DONE"),
        ("Readiness", readiness["status"], "DONE" if readiness["status"] == "ready" else "BLOCKED" if readiness["status"] == "not_ready" else "PENDING"),
    ]
    for label, value, status in checks:
        lines.append(f"- [{status}] {label}: {value}")
    return "\n".join(lines) + "\n"


def build_submission_memo(
    *,
    bundle: ReportBundle,
    customer_context: dict[str, Any],
    review_closure_status: dict[str, Any],
    readiness: dict[str, Any],
    included_customer_files: list[str],
) -> dict[str, Any]:
    meta = customer_context["document_meta"]
    unresolved = review_closure_status.get("unresolved_review_items", 0)
    return {
        "title": f"{meta.get('client_name', 'Client')} submission memo",
        "purpose": f"Deliver the final customer-facing assessment package for {meta.get('project_name', 'the engagement')}.",
        "included_artifacts": included_customer_files,
        "summary": [
            f"Assessment run `{bundle.run_id}` identified {bundle.summary.get('issues', 0)} customer-facing issues.",
            f"Release readiness is `{readiness['status']}`.",
            f"Unresolved review items at submission time: {unresolved}.",
        ],
        "notes": [
            "This package is intended for customer distribution and excludes internal-only analysis files.",
            "Tracker rows indicate whether each issue is included in the final delivery set.",
        ],
        "follow_up": [
            "Please confirm receipt of the package and the planned remediation owner for each open issue.",
            "Schedule retesting after fixes are applied to validate closure.",
        ],
        "contact": meta.get("contact_email", ""),
        "footer_notice": meta.get("footer_notice", ""),
    }


def render_submission_memo_markdown(memo: dict[str, Any]) -> str:
    lines = [
        f"# {memo['title']}",
        "",
        "## Delivery Purpose",
        f"- {memo['purpose']}",
        "",
        "## Included Deliverables",
    ]
    lines.extend([f"- `{path}`" for path in memo["included_artifacts"]] or ["- No customer deliverables were selected."])
    lines.extend(["", "## Key Summary"])
    lines.extend([f"- {item}" for item in memo["summary"]])
    lines.extend(["", "## Notes"])
    lines.extend([f"- {item}" for item in memo["notes"]])
    lines.extend(["", "## Requested Follow-Up"])
    lines.extend([f"- {item}" for item in memo["follow_up"]])
    lines.extend(["", "## Contact", f"- {memo['contact'] or 'N/A'}"])
    if memo.get("footer_notice"):
        lines.extend(["", "## Notice", f"- {memo['footer_notice']}"])
    return "\n".join(lines) + "\n"


def build_real_data_rehearsal_summary(
    *,
    bundle: ReportBundle,
    review_queue: list[dict[str, Any]],
    readiness: dict[str, Any],
    real_input_selection: dict[str, Any] | None,
) -> dict[str, Any]:
    selected_inputs: list[str] = []
    notes: list[str] = []
    manual_inputs: dict[str, Any] = {}
    status = "completed"
    if real_input_selection:
        selected_inputs = list(real_input_selection.get("selected_run_inputs", []))
        notes.extend(real_input_selection.get("notes", []))
        manual_inputs = dict(real_input_selection.get("manual_resolution", {}))
        if real_input_selection.get("status") != "selected":
            status = "incomplete"
            notes.append("Real-data rehearsal was not completed with automatically selected real inputs.")
        if any(item.get("manual_source") != "real_explicit" for item in manual_inputs.values()):
            status = "incomplete"
            notes.append("Manual support files were not all resolved from explicit real/manual inputs.")
    else:
        status = "incomplete"
        notes.append("Real input auto-selection was not used for this run.")
    return {
        "status": status,
        "run_id": bundle.run_id,
        "selected_inputs": selected_inputs,
        "findings_count": bundle.summary.get("deduped_findings", 0),
        "issues_count": bundle.summary.get("issues", 0),
        "review_queue_count": len(review_queue),
        "readiness_status": readiness.get("status"),
        "final_delivery_possible": readiness.get("status") == "ready",
        "manual_inputs": manual_inputs,
        "notes": notes,
    }


def render_real_data_rehearsal_summary(summary: dict[str, Any]) -> str:
    lines = [
        "# Real Data Rehearsal Summary",
        "",
        f"- status: `{summary['status']}`",
        f"- run_id: `{summary['run_id']}`",
        f"- findings_count: {summary['findings_count']}",
        f"- issues_count: {summary['issues_count']}",
        f"- review_queue_count: {summary['review_queue_count']}",
        f"- readiness_status: `{summary['readiness_status']}`",
        f"- final_delivery_possible: `{str(summary['final_delivery_possible']).lower()}`",
        "",
        "## Selected Real Inputs",
    ]
    lines.extend([f"- `{path}`" for path in summary["selected_inputs"]] or ["- No real input files were selected."])
    lines.extend(["", "## Manual Input Sources"])
    if summary.get("manual_inputs"):
        lines.extend(
            [
                f"- {key}: manual_source=`{value.get('manual_source', 'unclassified')}` "
                f"effective_path=`{value.get('effective_path')}`"
                for key, value in summary["manual_inputs"].items()
            ]
        )
    else:
        lines.append("- None.")
    lines.extend(["", "## Notes"])
    lines.extend([f"- {note}" for note in summary["notes"]] or ["- None."])
    return "\n".join(lines) + "\n"


def _closeout_summary(review_closure_status: dict[str, Any]) -> dict[str, int]:
    return {
        "resolved_review_count": review_closure_status.get("resolved_review_items", 0),
        "accepted_risk_count": review_closure_status.get("accepted_risk_items", 0),
        "remaining_blocker_count": review_closure_status.get("unresolved_review_items", 0),
    }


def _comparison_status_map(comparison: dict[str, Any]) -> dict[str, str]:
    mapping: dict[str, str] = {}
    for key in ["new_issues", "resolved_issues", "changed_issues", "unchanged_issues"]:
        for row in comparison.get(key, []):
            if row.get("cluster_key"):
                mapping[row["cluster_key"]] = key.removesuffix("_issues")
    return mapping


def _first(items: list[str], fallback: str) -> str:
    for item in items:
        if item:
            return item
    return fallback


def _tracker_columns() -> list[str]:
    return [
        "issue_id",
        "title",
        "severity",
        "confidence",
        "weakness_family",
        "primary_cwe",
        "affected_assets",
        "immediate_action",
        "structural_fix",
        "validation_after_fix",
        "recommended_owner",
        "target_due",
        "status",
        "source_run_id",
        "compared_run_status",
        "suppression_status",
        "review_required",
        "review_disposition",
        "approval_status",
        "final_delivery_included",
        "closeout_status",
    ]


def _dedupe_paths(paths: list[str]) -> list[str]:
    ordered: list[str] = []
    seen: set[str] = set()
    for path in paths:
        if path not in seen:
            ordered.append(path)
            seen.add(path)
    return ordered


def _base_document(title: str) -> Document:
    document = Document()
    section = document.sections[0]
    section.top_margin = Inches(0.8)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.8)
    section.right_margin = Inches(0.8)
    for style_name in ["Normal", "Heading 1", "Heading 2", "Heading 3"]:
        style = document.styles[style_name]
        style.font.name = "Malgun Gothic"
        style._element.rPr.rFonts.set(qn("w:eastAsia"), "Malgun Gothic")
    document.styles["Normal"].font.size = Pt(10.5)
    document.core_properties.title = title
    return document
